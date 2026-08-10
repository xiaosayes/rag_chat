"""
多格式文档加载器
支持 PDF、Word (docx)、TXT、Markdown、图片 (OCR)、PPT 等格式
统一输出为 Artifact 对象，方便后续入库检索
"""

from __future__ import annotations

import re
from pathlib import Path
from typing import Any, Dict, List, Optional
from dataclasses import dataclass, field

from loguru import logger

from src.data_loader import Artifact, DataLoader
from src.utils import generate_id

# bug-117b：文档加载统一清洗 C0/DEL 控制字符（PDF/Office 提取常带 \x00-\x1f 杂字符，
# 会污染切片与检索）。保留有含义的 \n \t \r。
_CONTROL_CHAR_RE = re.compile(r"[\x00-\x08\x0b\x0c\x0e-\x1f\x7f]")


# ========== 支持的文档类型 ==========

SUPPORTED_EXTENSIONS = {
    ".txt": "text/plain",
    ".md": "text/markdown",
    ".json": "application/json",
    ".csv": "text/csv",
    ".pdf": "application/pdf",
    ".docx": "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
    ".doc": "application/msword",
    ".pptx": "application/vnd.openxmlformats-officedocument.presentationml.presentation",
    ".ppt": "application/vnd.ms-powerpoint",
    ".png": "image/png",
    ".jpg": "image/jpeg",
    ".jpeg": "image/jpeg",
    ".gif": "image/gif",
    ".webp": "image/webp",
    ".bmp": "image/bmp",
    ".tiff": "image/tiff",
    ".tif": "image/tiff",
    ".xlsx": "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
}


@dataclass
class Document:
    """文档对象，表示一个文档文件及其内容"""
    path: Path
    content: str = ""
    title: str = ""
    metadata: Dict[str, Any] = field(default_factory=dict)
    pages: List[str] = field(default_factory=list)
    format: str = ""

    def __post_init__(self):
        if not self.title:
            self.title = self.path.stem
        if not self.format:
            self.format = self.path.suffix.lower()


class DocumentParser:
    """文档解析器基类"""

    extensions: List[str] = []

    def parse(self, path: Path) -> Document:
        raise NotImplementedError


class TxtParser(DocumentParser):
    """纯文本解析器"""

    extensions = [".txt", ".md", ".csv"]

    def parse(self, path: Path) -> Document:
        with open(path, "r", encoding="utf-8", errors="replace") as f:
            content = f.read()
        return Document(
            path=path,
            content=content,
            title=path.stem,
            format=path.suffix.lower(),
            metadata={"size": len(content), "lines": content.count("\n") + 1},
        )


class JsonParser(DocumentParser):
    """JSON 解析器（直接返回内容作为字符串）"""

    extensions = [".json"]

    def parse(self, path: Path) -> Document:
        import json
        with open(path, "r", encoding="utf-8") as f:
            data = json.load(f)
        content = json.dumps(data, ensure_ascii=False, indent=2)
        return Document(
            path=path,
            content=content,
            title=path.stem,
            format=".json",
            metadata={"size": len(content), "records": len(data) if isinstance(data, list) else 1},
        )


class PdfParser(DocumentParser):
    """PDF 解析器（使用 pypdf）"""

    extensions = [".pdf"]

    def parse(self, path: Path) -> Document:
        try:
            from pypdf import PdfReader
        except ImportError:
            raise ImportError("请先安装 pypdf: pip install pypdf")

        reader = PdfReader(str(path))
        pages = []
        for i, page in enumerate(reader.pages):
            text = page.extract_text() or ""
            pages.append(text)

        content = "\n\n".join(pages)
        metadata = {
            "pages": len(pages),
            "size": len(content),
            "pdf_version": str(reader.pdf_header) if hasattr(reader, 'pdf_header') else "",
        }
        # 尝试从元数据获取标题
        title = path.stem
        if reader.metadata and reader.metadata.title:
            title = reader.metadata.title

        return Document(
            path=path,
            content=content,
            title=title,
            format=".pdf",
            pages=pages,
            metadata=metadata,
        )


class DocxParser(DocumentParser):
    """Word 文档解析器（支持 .docx 和 .doc）"""

    extensions = [".docx", ".doc"]

    def parse(self, path: Path) -> Document:
        ext = path.suffix.lower()
        if ext == ".doc":
            logger.warning(
                f"{path.name} 是旧版 .doc 格式，python-docx 不支持。"
                f"建议用 Word 另存为 .docx 格式。将尝试按纯文本读取。"
            )
            # 尝试作为纯文本读取（至少能提取部分内容）
            try:
                with open(path, "r", encoding="utf-8", errors="replace") as f:
                    content = f.read()
                return Document(
                    path=path,
                    content=content,
                    title=path.stem,
                    format=".doc",
                    metadata={"size": len(content), "warning": "旧版 .doc 格式，内容可能不完整"},
                )
            except Exception:
                return Document(
                    path=path,
                    content=f"[无法解析 .doc 格式: {path.name}]",
                    title=path.stem,
                    format=".doc",
                    metadata={"error": "python-docx 不支持旧版 .doc 格式"},
                )

        try:
            from docx import Document as DocxDocument
        except ImportError:
            raise ImportError("请先安装 python-docx: pip install python-docx")

        doc = DocxDocument(str(path))
        paragraphs = []
        for para in doc.paragraphs:
            if para.text.strip():
                paragraphs.append(para.text)

        # 提取表格内容
        tables_text = []
        for table in doc.tables:
            for row in table.rows:
                cells = [cell.text.strip() for cell in row.cells]
                tables_text.append(" | ".join(cells))

        content = "\n".join(paragraphs)
        if tables_text:
            content += "\n\n[表格内容]\n" + "\n".join(tables_text)

        metadata = {
            "paragraphs": len(paragraphs),
            "tables": len(doc.tables),
            "size": len(content),
        }

        return Document(
            path=path,
            content=content,
            title=path.stem,
            format=".docx",
            metadata=metadata,
        )


class PptxParser(DocumentParser):
    """PPT 解析器"""

    extensions = [".pptx", ".ppt"]

    def parse(self, path: Path) -> Document:
        # audit-F20 修复：.ppt 为旧版二进制格式，python-pptx 不支持（此前声称支持
        # 但必然抛错）。与 DocxParser 对 .doc 的处理一致：告警 + 纯文本兜底。
        if path.suffix.lower() == ".ppt":
            logger.warning(
                f"{path.name} 是旧版 .ppt 格式，python-pptx 不支持。"
                f"建议用 PowerPoint 另存为 .pptx 格式。将尝试按纯文本读取。"
            )
            try:
                with open(path, "r", encoding="utf-8", errors="replace") as f:
                    content = f.read()
                return Document(
                    path=path,
                    content=content,
                    title=path.stem,
                    format=".ppt",
                    metadata={"size": len(content), "warning": "旧版 .ppt 格式，内容可能不完整"},
                )
            except Exception:
                return Document(
                    path=path,
                    content=f"[无法解析 .ppt 格式: {path.name}]",
                    title=path.stem,
                    format=".ppt",
                    metadata={"error": "python-pptx 不支持旧版 .ppt 格式"},
                )
        try:
            from pptx import Presentation
        except ImportError:
            raise ImportError("请先安装 python-pptx: pip install python-pptx")

        prs = Presentation(str(path))
        slides_text = []
        for i, slide in enumerate(prs.slides, 1):
            slide_content = [f"--- 幻灯片 {i} ---"]
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_content.append(shape.text)
                if shape.has_table:
                    table = shape.table
                    for row in table.rows:
                        cells = [cell.text.strip() for cell in row.cells]
                        slide_content.append(" | ".join(cells))
            slides_text.append("\n".join(slide_content))

        content = "\n\n".join(slides_text)
        metadata = {
            "slides": len(prs.slides),
            "size": len(content),
        }

        return Document(
            path=path,
            content=content,
            title=path.stem,
            format=".pptx",
            metadata=metadata,
        )


class ImageParser(DocumentParser):
    """
    图片解析器（OCR）
    优先使用 PaddleOCR（GPU 加速，中文识别效果最好）
    备选方案：Tesseract OCR
    """

    extensions = [".png", ".jpg", ".jpeg", ".gif", ".webp", ".bmp", ".tiff", ".tif"]

    def __init__(self, lang: str = "ch", use_paddle: bool = True):
        self.lang = lang
        self.use_paddle = use_paddle
        self._paddle_ocr = None
        self._pytesseract_available = None

    def parse(self, path: Path) -> Document:
        """对图片进行 OCR 识别"""
        if self.use_paddle:
            try:
                return self._parse_with_paddleocr(path)
            except Exception as e:
                logger.warning(f"PaddleOCR 识别失败，尝试 Tesseract: {e}")

        return self._parse_with_tesseract(path)

    def _parse_with_paddleocr(self, path: Path) -> Document:
        """使用 PaddleOCR 进行文字识别（自动检测 GPU 可用性）"""
        if self._paddle_ocr is None:
            try:
                from paddleocr import PaddleOCR
                # 检测 GPU 是否可用
                use_gpu = False
                try:
                    import paddle
                    use_gpu = paddle.is_compiled_with_cuda()
                except ImportError:
                    pass
                self._paddle_ocr = PaddleOCR(
                    use_angle_cls=True,
                    lang=self.lang,
                    use_gpu=use_gpu,
                    show_log=False,
                )
                mode_name = "GPU" if use_gpu else "CPU"
                logger.info(f"PaddleOCR 初始化完成（{mode_name} 模式）")
            except ImportError:
                raise ImportError("请先安装 PaddleOCR: pip install paddleocr paddlepaddle-gpu")

        result = self._paddle_ocr.ocr(str(path), cls=True)
        if not result or not result[0]:
            return Document(
                path=path,
                content="[图片未识别出文字]",
                title=path.stem,
                format=path.suffix.lower(),
                metadata={"ocr_engine": "paddleocr", "text_found": False},
            )

        # 提取文本（bug-058 修复：兼容 PaddleOCR 2.x 与 3.x 输出格式）
        lines = []
        for line in result[0]:
            if not isinstance(line, (list, tuple)) or len(line) < 2:
                continue
            # 2.x: [box, (text, confidence)]；3.x: [text, confidence]
            if isinstance(line[1], (list, tuple)):
                text = line[1][0]
                confidence = line[1][1]
            else:
                text = line[0]
                confidence = line[1]
            if confidence > 0.3:  # 低置信度过滤
                lines.append(text)

        content = "\n".join(lines)
        return Document(
            path=path,
            content=content,
            title=path.stem,
            format=path.suffix.lower(),
            metadata={
                "ocr_engine": "paddleocr",
                "text_lines": len(lines),
                "text_found": len(lines) > 0,
            },
        )

    def _parse_with_tesseract(self, path: Path) -> Document:
        """使用 Tesseract OCR（备选方案）"""
        try:
            import pytesseract
            from PIL import Image
        except ImportError:
            raise ImportError("请先安装 pytesseract 和 Pillow: pip install pytesseract Pillow")

        # 检查 Tesseract 是否安装
        if self._pytesseract_available is None:
            try:
                pytesseract.get_tesseract_version()
                self._pytesseract_available = True
            except Exception:
                self._pytesseract_available = False
                raise RuntimeError(
                    "Tesseract OCR 未安装。\n"
                    "Linux: sudo apt-get install tesseract-ocr tesseract-ocr-chi-sim\n"
                    "Windows: 下载安装 https://github.com/UB-Mannheim/tesseract/wiki"
                )

        image = Image.open(str(path))
        lang_code = "chi_sim+eng" if self.lang == "ch" else "eng"
        content = pytesseract.image_to_string(image, lang=lang_code)

        return Document(
            path=path,
            content=content,
            title=path.stem,
            format=path.suffix.lower(),
            metadata={
                "ocr_engine": "tesseract",
                "text_lines": content.count("\n"),
                "text_found": len(content.strip()) > 0,
            },
        )


# ========== 文档加载器（统一入口） ==========

class DocumentLoader:
    """
    统一文档加载器
    自动检测文件格式并选择合适的解析器
    """

    def __init__(self, enable_ocr: bool = True, ocr_engine: str = "paddle"):
        self.enable_ocr = enable_ocr
        self.ocr_engine = ocr_engine

        # 注册解析器
        self.parsers: List[DocumentParser] = [
            TxtParser(),
            JsonParser(),
            PdfParser(),
            DocxParser(),
            PptxParser(),
        ]
        if enable_ocr:
            self.parsers.append(
                ImageParser(use_paddle=(ocr_engine == "paddle"))
            )

        # 构建扩展名 -> 解析器映射
        self.extension_map: Dict[str, DocumentParser] = {}
        for parser in self.parsers:
            for ext in parser.extensions:
                self.extension_map[ext.lower()] = parser

    def load_file(self, path: Path) -> Document:
        """加载单个文件"""
        # bug-023 修复：解析真实路径（规范化符号链接/相对路径）
        # 注意：本方法仅做路径规范化与存在性检查，不做目录归属校验；
        # 若未来暴露给外部输入（如 Web 上传路径），需追加 containment 校验（audit-F26）
        try:
            resolved = path.resolve()
        except (OSError, RuntimeError) as e:
            raise ValueError(f"无法解析路径: {path} - {e}")
        if not resolved.exists():
            raise FileNotFoundError(f"文件不存在: {resolved}")

        ext = path.suffix.lower()
        if ext not in self.extension_map:
            raise ValueError(
                f"不支持的文件格式: {ext}\n"
                f"支持的格式: {', '.join(SUPPORTED_EXTENSIONS.keys())}"
            )

        parser = self.extension_map[ext]
        logger.info(f"正在解析: {path.name} (格式: {ext})")
        doc = parser.parse(path)
        # bug-117b：统一清洗控制字符（对所有解析器生效，含 PDF/Office 提取的杂字符）
        doc.content = _CONTROL_CHAR_RE.sub("", doc.content)
        if doc.pages:
            doc.pages = [_CONTROL_CHAR_RE.sub("", p) for p in doc.pages]
        logger.info(f"解析完成: {path.name} → {len(doc.content)} 字符")
        return doc

    def load_directory(
        self,
        directory: Path,
        recursive: bool = True,
        extensions: Optional[List[str]] = None,
    ) -> List[Document]:
        """加载目录下的所有支持文件"""
        if not directory.exists():
            raise FileNotFoundError(f"目录不存在: {directory}")

        documents = []
        pattern = "**/*" if recursive else "*"

        for file_path in sorted(directory.glob(pattern)):
            if not file_path.is_file():
                continue
            ext = file_path.suffix.lower()
            if extensions and ext not in extensions:
                continue
            if ext not in self.extension_map:
                continue

            try:
                doc = self.load_file(file_path)
                documents.append(doc)
            except Exception as e:
                logger.error(f"解析失败 {file_path.name}: {e}")
                continue

        logger.info(f"目录加载完成: {directory} → {len(documents)} 个文档")
        return documents

    def document_to_artifact(self, doc: Document, category: str = "文档资料") -> Artifact:
        """将 Document 转换为 Artifact 对象"""
        # 从内容中提取标题和描述
        content = doc.content
        title = doc.title or doc.path.stem

        # 截取前 5000 字作为描述（bug-011：原 500 字截断导致长文档大量信息丢失）
        description = content[:5000] if content else ""

        # 提取标签
        tags = [category, doc.format.lstrip(".")]
        if doc.metadata.get("text_found") is False:
            tags.append("未识别文字")

        return Artifact(
            name=title,
            dynasty="",
            category=category,
            material="",
            location="",
            description=description,
            historical_significance=f"来自文档: {doc.path.name}",
            cultural_value=f"格式: {doc.format}，大小: {len(content)} 字符",
            tags=tags,
            importance=3,
            extra={
                "source_file": str(doc.path),
                "source_format": doc.format,
                "full_content": content,
                "metadata": doc.metadata,
            },
        )

    def load_all_as_artifacts(
        self,
        source: Path,
        category: str = "文档资料",
        recursive: bool = True,
    ) -> List[Artifact]:
        """
        从文件或目录加载所有文档，统一转换为 Artifact 对象
        方便后续切片和入库
        """
        artifacts: List[Artifact] = []
        if source.is_file():
            # Excel 表格：每行一条记录，委托 DataLoader 解析（不经过单文档模型，bug-109）
            if source.suffix.lower() == ".xlsx":
                return DataLoader.load(source)
            docs = [self.load_file(source)]
        elif source.is_dir():
            docs = self.load_directory(source, recursive=recursive)
            # Excel 表格单独收集（load_directory 的 extension_map 不含 .xlsx，
            # 且 Excel 是多记录文件，不走单文档模型）
            pattern = "**/*.xlsx" if recursive else "*.xlsx"
            for xlsx_path in sorted(source.glob(pattern)):
                if xlsx_path.is_file():
                    try:
                        artifacts.extend(DataLoader.load(xlsx_path))
                    except Exception as e:
                        logger.error(f"解析失败 {xlsx_path.name}: {e}")
        else:
            raise FileNotFoundError(f"路径不存在: {source}")

        for doc in docs:
            # P1-5 修复：长文档按段切分为多个 Artifact，
            # 避免 document_to_artifact 的 5000 字符截断导致后续内容完全无法被检索
            content = doc.content
            if len(content) <= 5000:
                artifacts.append(self.document_to_artifact(doc, category=category))
                continue
            title = doc.title or doc.path.stem
            segment_size = 4500
            total_segments = (len(content) + segment_size - 1) // segment_size
            for idx, start in enumerate(range(0, len(content), segment_size)):
                seg_doc = Document(
                    path=doc.path,
                    content=content[start:start + segment_size],
                    title=f"{title}（第{idx + 1}/{total_segments}部分）",
                    metadata={**doc.metadata, "segment": idx + 1, "segment_total": total_segments},
                    format=doc.format,
                )
                artifacts.append(self.document_to_artifact(seg_doc, category=category))
            logger.info(
                f"长文档 {doc.path.name} 已切分为 {total_segments} 段（总 {len(content)} 字符）"
            )

        logger.info(f"文档转换完成: {len(docs)} 个文档 → {len(artifacts)} 个 Artifact")
        return artifacts


# ========== 便捷函数 ==========

def load_document(file_path: str) -> Document:
    """便捷函数：加载单个文档"""
    loader = DocumentLoader()
    return loader.load_file(Path(file_path))


def load_documents_to_artifacts(
    source_path: str,
    category: str = "文档资料",
) -> List[Artifact]:
    """便捷函数：加载文档并转换为 Artifact"""
    loader = DocumentLoader()
    return loader.load_all_as_artifacts(Path(source_path), category=category)